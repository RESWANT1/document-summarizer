"""
extractors/pptx_extractor.py
Extracts slide text from PowerPoint (.pptx) files using python-pptx.
Iterates over slides, shapes, and text frames in order.
"""

import logging
from pptx import Presentation

logger = logging.getLogger(__name__)


def extract_pptx(filepath: str) -> str:
    """Return concatenated text from all slides of a .pptx file."""
    prs = Presentation(filepath)
    parts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        slide_texts.append(txt)
        if slide_texts:
            parts.append(f"[Slide {slide_num}] " + " ".join(slide_texts))

    full_text = "\n".join(parts)
    logger.info("PPTX extraction complete: %d slides, %d characters.", len(prs.slides), len(full_text))
    return full_text